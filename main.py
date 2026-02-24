from pptx import Presentation
from collections import Counter
import sys
import os

def extract_slides(filepath):
    prs = Presentation(filepath)
    slides_data = []
    for i, slide in enumerate(prs.slides, start=1):
        title = "No title"
        body_text = []
        has_images = False
        has_tables = False
        notes = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title = shape.text.strip()
            if shape.shape_type == 13:
                has_images = True
            if shape.has_table:
                has_tables = True
            elif shape.has_text_frame and not (shape.is_placeholder and shape.placeholder_format.idx == 0):
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        body_text.append(text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides_data.append({
            "number": i,
            "title": title,
            "body_text": body_text,
            "has_images": has_images,
            "has_tables": has_tables,
            "notes": notes
        })
    return slides_data

def generate_synopsis(slides_data):
    print("=" * 50)
    print("  PRESENTATION SYNOPSIS")
    print("=" * 50)
    print(f"  Total slides: {len(slides_data)}")
    print("=" * 50)
    print("\nSLIDE BREAKDOWN")
    print("-" * 30)
    for slide in slides_data:
        print(f"\n  Slide {slide['number']}: {slide['title']}")
        if slide["has_images"]:
            print("    [Contains images]")
        if slide["has_tables"]:
            print("    [Contains tables]")
        if slide["notes"]:
            note = slide["notes"]
            if len(note) > 100:
                note = note[:100] + "..."
            print(f"    Notes: {note}")
        for line in slide["body_text"]:
            print(f"    - {line}")
    print("\nKEY TOPICS")
    print("_" * 30)
    all_words = []
    for slide in slides_data:
        for line in slide["body_text"]:
            words = line.lower().split()
            all_words.extend(words)
    stop_words = {"the", "a", "an", "is", "are", "was", "were", "to", "of",
                  "in", "for", "on", "with", "at", "by", "from", "and", "or",
                  "but", "not", "you", "your", "it", "its", "this", "that",
                  "be", "as", "if", "do", "does", "has", "have", "will",
                  "can", "should", "may", "about", "what", "which", "who",
                  "how", "when", "where", "than", "more", "some", "any",
                  "other", "into", "no", "so", "up", "out", "just", "also"}
    cleaned = [w.strip(".,!?;:()[]\"'-") for w in all_words if w.strip(".,!?;:()[]\"'-") not in stop_words and len(w.strip(".,!?;:()[]\"'-")) > 2]
    counts = Counter(cleaned)
    for word, count in counts.most_common(15):
        bar = "#" * count
        print(f"    {word:<20} {bar} ({count})")

def group_slides_by_topic(slides_data):
    groups = []
    current_group = None

    for slide in slides_data:
        title = slide["title"]
        if current_group is None or title != current_group["title"]:
            current_group = {"title": title, "slides": [slide]}
            groups.append(current_group)
        else:
            current_group["slides"].append(slide)

    return groups

def write_synopsis(slides_data):
    groups = group_slides_by_topic(slides_data)
    print("SYNOPSIS")
    print("=" * 50)
    print()
    first_body = slides_data[0]["body_text"]
    topic = first_body[0] if first_body else slides_data[0]["title"].split("\n")[0]
    print(f"This presentation on \"{topic}\" covers {len(slides_data)} slides across the following topics.")
    section_titles = []
    for i, group in enumerate(groups):
        clean_title = group["title"].split("\n")[0].strip()
        if i == 0 and len(groups[0]["slides"][0]["body_text"]) <= 1:
            continue
        if "activity" in clean_title.lower() or "thank" in clean_title.lower():
            continue
        if clean_title not in section_titles:
            section_titles.append(clean_title)
    print(f"The main topics covered are: {', '.join(section_titles[:6])}.")
    print()
    transitions = ["The presentation begins by discussing", "It then covers", "Next, it explores", "Following this,", "The presentation also addresses", "Additionally, it examines", "It further discusses", "The presentation goes on to explain", "Moreover,", "Finally,"]
    for idx, group in enumerate(groups):
        title = group["title"].split("\n")[0].strip()
        slides = group["slides"]
        if idx == 0 and len(groups[0]["slides"][0]["body_text"]) <= 1:
            continue
        if "activity" in title.lower() or "thank" in title.lower():
            continue
        all_text = []
        for slide in slides:
            all_text.extend(slide["body_text"])
        if not all_text:
            continue
        transition = transitions[idx % len(transitions)]
        if all(len(t) < 80 for  t in all_text[:3]):
            summary = ", ".join(t.rstrip(":;,") for t in all_text[:3])
        else:
            summary = " ".join(all_text[:3])
        if len(summary) > 300:
            summary = summary[:297] + "..."
        print(f"{transition} {title.lower()}. {summary}")
    image_count = sum(1 for s in slides_data if s["has_images"])
    notes_count = sum(1 for s in slides_data if s["notes"])
    print(f"The presentation uses {image_count} images across its slides", end="")
    if notes_count > 0:
        print(f" and includes speaker notes on {notes_count} slides.", end="")
    print()



filepath = "Research Skills for Academic Writing.pptx"
if not os.path.isfile(filepath):
    print(f"Error: File not found: {filepath}")
    sys.exit(1)
if not filepath.lower().endswith(".pptx"):
    print(f"Error: Expected a .pptx file, got: {filepath}")
    sys.exit(1)

data = extract_slides(filepath)
import io
import contextlib

output_summary = filepath.replace(".pptx", "_synopsis.txt")
with open(output_summary, "w", encoding="utf-8") as f:
    buffer = io.StringIO()
    with contextlib.redirect_stdout(buffer):
        write_synopsis(data)
    f.write(buffer.getvalue())

output_detailed = filepath.replace(".pptx", "_detailed.txt")
with open(output_detailed, "w", encoding="utf-8") as f:
    buffer = io.StringIO()
    with contextlib.redirect_stdout(buffer):
        generate_synopsis(data)
        print("\n")
        write_synopsis(data)
    f.write(buffer.getvalue())

write_synopsis(data)
print(f"\nSummary saved to: {output_summary}")
print(f"Detailed saved to: {output_detailed}")